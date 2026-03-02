#!/usr/bin/env python3
"""
Generate a professional PowerPoint presentation from report.md
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Color scheme - Executive style
COLOR_PRIMARY = RGBColor(25, 45, 85)      # Deep Navy
COLOR_ACCENT = RGBColor(0, 120, 180)      # Professional Blue
COLOR_LIGHT = RGBColor(240, 245, 250)     # Light Blue
COLOR_TEXT = RGBColor(40, 40, 40)         # Dark Gray
COLOR_LIGHT_TEXT = RGBColor(100, 100, 100) # Medium Gray
COLOR_WHITE = RGBColor(255, 255, 255)

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Define slide layouts
    blank_slide_layout = prs.slide_layouts[6]  # Blank layout

    # Slide 1: Title Slide
    slide1 = prs.slides.add_slide(blank_slide_layout)
    add_background(slide1, COLOR_PRIMARY)

    # Title
    title_box = slide1.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(2))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = "Strategic Analysis: Agentic AI Integration"
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE

    # Subtitle
    subtitle_box = slide1.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(9), Inches(1.5))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.word_wrap = True
    p = subtitle_frame.paragraphs[0]
    p.text = "Claude Code vs. Google Antigravity vs. GitHub Copilot"
    p.font.size = Pt(28)
    p.font.color.rgb = COLOR_ACCENT

    # Date
    date_box = slide1.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(0.6))
    date_frame = date_box.text_frame
    p = date_frame.paragraphs[0]
    p.text = "March 2, 2026"
    p.font.size = Pt(18)
    p.font.color.rgb = COLOR_LIGHT

    # Slide 2: Executive Summary
    slide2 = add_section_slide(prs, "Executive Summary")

    content = [
        "Strategic Objective:",
        "• Collect, enhance, and organize a robust library of open-source and self-developed skills in ./claude/skills/",
        "",
        "Key Finding:",
        "• Claude Code, Google Antigravity, and GitHub Copilot all support skills invocation",
        "• Only Claude Code supports reusable subagents in ./claude/agents/",
        "",
        "Critical Advantage:",
        "• Subagents are essential for operationalizing complex engineering workflows across the team",
        "• Enables hierarchical workflows, parallel execution, and built-in work summarization"
    ]

    add_content_to_slide(slide2, content)

    # Slide 3: Feature Comparison - Claude Code
    slide3 = add_section_slide(prs, "Feature Comparison: Claude Code")

    content = [
        "Custom Subagents: ✓ YES",
        "• Define specialized agents in ./claude/agents/ with full framework support",
        "• Each subagent operates with isolated context and can run in parallel",
        "• Supports worktree isolation for git-safe concurrent execution",
        "• Agents can delegate to other agents, creating hierarchical workflows",
        "• Built-in work summarization back to main conversation",
        "",
        "Skills Invocation: ✓ YES",
        "• Automatically triggers skills based on conversation context",
        "• Skills in ./claude/skills/ available globally to main agent and all subagents",
        "",
        "Extensibility: COMPREHENSIVE",
        "• CLAUDE.md for persistent team context and standards",
        "• settings.json for event-driven hooks (automation)",
        "• Full MCP server integration",
        "",
        "Pricing: $17–20/month individual | $25/seat/month team"
    ]

    add_content_to_slide(slide3, content)

    # Slide 4: Feature Comparison - Google Antigravity
    slide4 = add_section_slide(prs, "Feature Comparison: Google Antigravity")

    content = [
        "Custom Subagents: ✗ NO",
        "• No equivalent to ./claude/agents/ framework",
        "• Uses implicit 'Manager' views (Task, Plan, Code Review) but cannot define custom agents",
        "• Complex orchestration must be guided manually",
        "",
        "Skills Invocation: ✓ YES",
        "• Supports invoking skills from ./claude/skills/",
        "• Recent addition (January 2026), ecosystem still maturing",
        "",
        "Extensibility: MODERATE",
        "• Unique strength in artifact-based transparency",
        "• Lacks native hooks system or CLAUDE.md equivalent",
        "",
        "Pricing: Free (preview) | ~$20/month individual (expected later in 2026)"
    ]

    add_content_to_slide(slide4, content)

    # Slide 5: Feature Comparison - GitHub Copilot
    slide5 = add_section_slide(prs, "Feature Comparison: GitHub Copilot")

    content = [
        "Custom Subagents: ✗ NO",
        "• Supports integrated skills but lacks user-defined subagent feature",
        "• Cannot define specialized, reusable agents for specific workflow parts",
        "",
        "Skills Invocation: ✓ YES",
        "• Supports 'Agent Skills' (folders of instructions and scripts)",
        "• Can use skills in .github/skills/ or .claude/skills/",
        "",
        "Extensibility: MODERATE",
        "• Integrates with MCP servers",
        "• Supports custom instructions via instructions.md",
        "• No native parallel subagent delegation framework",
        "",
        "Pricing: $10/month (Pro) | $39/month (Pro+)"
    ]

    add_content_to_slide(slide5, content)

    # Slide 6: Competitive Matrix
    slide6 = add_section_slide(prs, "Competitive Matrix")

    # Add table
    rows, cols = 8, 4
    left = Inches(0.4)
    top = Inches(1.3)
    width = Inches(9.2)
    height = Inches(5.5)

    table_shape = slide6.shapes.add_table(rows, cols, left, top, width, height).table

    # Set column widths
    table_shape.columns[0].width = Inches(3.5)
    table_shape.columns[1].width = Inches(2)
    table_shape.columns[2].width = Inches(2)
    table_shape.columns[3].width = Inches(2)

    # Headers
    headers = ["Capability", "Claude Code", "Antigravity", "GitHub Copilot"]
    for col, header in enumerate(headers):
        cell = table_shape.cell(0, col)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_PRIMARY
        tf = cell.text_frame
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = COLOR_WHITE
        tf.paragraphs[0].font.size = Pt(11)

    # Data
    data = [
        ["Skill Integration (./claude/skills/)", "✅", "✅", "✅"],
        ["Custom Subagents (./claude/agents/)", "✅", "❌", "❌"],
        ["Team-wide Agent Sharing", "✅", "❌ Manual", "❌ Manual"],
        ["Parallel Execution", "✅", "⚠️ Implicit", "❌"],
        ["Hooks & Automation", "✅", "❌", "⚠️ Limited"],
        ["Artifact Transparency", "❌ Chat-only", "✅", "❌ Chat-only"],
        ["Starting Cost", "$17/month", "Free", "$10/month"]
    ]

    for row, row_data in enumerate(data, start=1):
        for col, cell_text in enumerate(row_data):
            cell = table_shape.cell(row, col)
            cell.text = cell_text
            if row % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLOR_LIGHT
            tf = cell.text_frame
            tf.paragraphs[0].font.size = Pt(10)
            if col == 0:
                tf.paragraphs[0].font.bold = True
                tf.paragraphs[0].font.color.rgb = COLOR_PRIMARY
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Slide 7: Strategic Justification
    slide7 = add_section_slide(prs, "Strategic Justification")

    content = [
        "The Central Insight:",
        "As skills are collected and organized, Claude Code is the only platform that allows leverage of",
        "that investment into reusable agents.",
        "",
        "Scenario 1: With Claude Code",
        "• Define 'React Architecture Specialist' agent once in ./claude/agents/",
        "• Every structural review invokes that specific agent automatically",
        "• Agent knows patterns, uses relevant skills, works independently",
        "• Build logic once, use forever",
        "",
        "Scenario 2: With Antigravity or GitHub Copilot",
        "• Re-explain context for each task",
        "• Manually re-select skills",
        "• Manually oversee every workflow",
        "• Duplicated effort, not automation",
        "",
        "Bottom Line:",
        "Neither Copilot nor Antigravity solve the core challenge: operationalizing engineering DNA",
        "into repeatable, autonomous workflows."
    ]

    add_content_to_slide(slide7, content)

    # Slide 8: Targeted Skills Capabilities
    slide8 = add_section_slide(prs, "Targeted Skills Capabilities")

    content = [
        "Planned Skill Domains:",
        "",
        "UI & Components:",
        "• shadcn-ui, ai-elements, building-components, web-design-guidelines",
        "",
        "React & Next.js:",
        "• Composition patterns, cache components, upgrade paths",
        "",
        "Backend & Auth:",
        "• Auth.js integration, MCP server management",
        "",
        "Video & Media:",
        "• Remotion-based production, streamdown workflows",
        "",
        "Automation & Testing:",
        "• Browser automation (Playwright), visual regression",
        "",
        "Outcome:",
        "Once organized, these skills represent internal best practices. With Claude Code, specialized",
        "agents apply practices consistently rather than relying on manual oversight."
    ]

    add_content_to_slide(slide8, content)

    # Slide 9: Rollout Plan Overview
    slide9 = add_section_slide(prs, "Proposed Rollout Plan")

    content = [
        "Four-Phase Implementation Strategy:",
        "",
        "Phase 1: Software Development Skills & Agents",
        "Core development infrastructure with specialized agents for full-stack engineering",
        "",
        "Phase 2: Quality Assurance & Testing Agents",
        "Comprehensive automated testing and quality validation infrastructure",
        "",
        "Phase 3: UI/UX Design Agents",
        "Design system governance and consistency through specialized agents",
        "",
        "Phase 4: Cybersecurity & Compliance Agents",
        "Security validation and compliance auditing capabilities",
        "",
        "Overall Success Metric:",
        "Automate 40% of routine engineering tasks through subagents within 90 days"
    ]

    add_content_to_slide(slide9, content)

    # Slide 10: Phase 1 - Software Development (Skills & Agents)
    slide10 = add_section_slide(prs, "Phase 1: Software Development Skills")

    content = [
        "Skills to Leverage:",
        "• Next.js Best Practices: File conventions, RSC boundaries, async patterns, metadata, error handling",
        "• Next.js Cache Components: PPR, use cache directive, cacheLife, cacheTag, updateTag optimization",
        "• Next.js Upgrade Paths: Automated migration guidance and codemods",
        "• React Composition Patterns: Compound components, render props, context providers",
        "• React Performance Best Practices: Memoization, code splitting, bundle optimization, hydration",
        "• Auth.js Integration: OAuth, credentials provider, environment configuration",
        "• MCP Server Skills: Building MCP servers with shared Zod schemas",
        "",
        "Skills to Build:",
        "• Django, FastAPI, Turborepo, Docker, Kubernetes, CI/CD Pipeline Automation",
        "• Automated Software Documentation, Release Management, Test Reporting"
    ]

    add_content_to_slide(slide10, content)

    # Slide 11: Phase 1 - Specialized Agents
    slide11 = add_section_slide(prs, "Phase 1: Specialized Development Agents")

    content = [
        "Specialized Agents:",
        "",
        "React Architecture Specialist Agent",
        "• Reviews component structure, enforces composition patterns, identifies prop drilling",
        "",
        "Next.js Optimization Agent",
        "• Audits configuration, suggests caching strategies, recommends PPR/RSC implementations",
        "",
        "Backend Architecture Agent",
        "• Validates API design, reviews database patterns, checks authentication flows",
        "",
        "DevOps Configuration Agent",
        "• Reviews Docker/Kubernetes manifests, validates CI/CD pipelines",
        "",
        "Upgrade & Migration Agent",
        "• Guides version upgrades, applies codemods, validates breaking changes",
        "",
        "Documentation & Release Agent",
        "• Generates product specs, technical documentation, release notes, manages changelogs"
    ]

    add_content_to_slide(slide11, content)

    # Slide 12: Phase 2 - QA & Testing
    slide12 = add_section_slide(prs, "Phase 2: Quality Assurance & Testing Agents")

    content = [
        "Skills to Leverage:",
        "• Playwright CLI: Browser automation, form filling, screenshots, page navigation, web testing",
        "• Browser Automation: Agent-driven interaction, data extraction, workflow testing",
        "• Before & After Comparison: Visual regression detection, screenshot documentation",
        "",
        "Skills to Build:",
        "• Web Application Testing, Python Testing & Connectors, Java Testing & Connectors",
        "• API Testing & Validation, Visual Regression Testing, Performance & Load Testing",
        "• Standardized Engineering Testing, Advanced Quality Assurance",
        "",
        "Specialized Agents:",
        "• Web Test Automation Agent • API Testing Agent • Visual Regression Agent",
        "• Python Connector Validator Agent • Java Connector Validator Agent",
        "• Performance Testing Agent • QA Automation & Resilience Agent",
        "",
        "Deliverables: Test automation library + operational QA agents + coverage & baseline repositories"
    ]

    add_content_to_slide(slide12, content)

    # Slide 13: Phase 3 - UI/UX Design
    slide13 = add_section_slide(prs, "Phase 3: UI/UX Design Agents")

    content = [
        "Skills to Leverage:",
        "• shadcn/ui, ai-elements, building-components, web-design-guidelines",
        "• Design System Documentation, JSON Render Core & React, Remotion & Remotion Best Practices",
        "• Streamdown: Markdown rendering, syntax highlighting, Mermaid diagrams, CJK support",
        "",
        "Skills to Build:",
        "• Accessibility Auditing, Design Tokens Management, Component Documentation",
        "• Brand & Theming, UX Writing & Microcopy",
        "",
        "Five Specialized Agents:",
        "• Component Library Specialist Agent • Design System Governance Agent",
        "• Accessibility Auditor Agent • UX Content Agent • Design-to-Code Agent",
        "",
        "Deliverables: Component & design library + 5 operational design agents + design tokens & documentation"
    ]

    add_content_to_slide(slide13, content)

    # Slide 14: Phase 4 - Cybersecurity & Compliance
    slide14 = add_section_slide(prs, "Phase 4: Cybersecurity & Compliance Agents")

    content = [
        "Skills to Leverage:",
        "• Playwright CLI & Browser Automation: Automated security testing, OWASP validation",
        "",
        "Skills to Build:",
        "• Penetration Testing Framework, Vulnerability Scanning, Compliance Auditing",
        "• Infrastructure Security, Security Code Review, API Security Testing",
        "• Incident Response & Threat Analysis",
        "",
        "Six Specialized Agents:",
        "• Security Code Review Agent • Penetration Testing Agent • Compliance Auditor Agent",
        "• Vulnerability Scanner Agent • Infrastructure Security Agent • API Security Tester Agent",
        "",
        "Deliverables: Security skills library + 6 operational agents + audit reports, testing playbooks,",
        "incident response procedures, and security baseline documentation"
    ]

    add_content_to_slide(slide14, content)

    # Slide 15: Summary of Findings
    slide15 = add_section_slide(prs, "Summary of Findings")

    content = [
        "Strategic Recommendation: ADOPT CLAUDE CODE",
        "",
        "Why Claude Code?",
        "• Aligns most closely with strategic goals for agentic workflow automation",
        "• Dedicated infrastructure for turning skill collection into autonomous, reusable agents",
        "• Only platform supporting custom subagent definition and parallel execution",
        "",
        "Alternative Platforms:",
        "• GitHub Copilot: Excellent for individual productivity, cost-effective ($10/month)",
        "• Google Antigravity: Excellent visual feedback, currently free (preview phase)",
        "• Neither solves the core challenge: operationalizing engineering DNA into repeatable",
        "  autonomous workflows",
        "",
        "Recommended Next Steps:",
        "• Pilot with a small group on the 'React Architecture Review' agent workflow",
        "• Validate platform utility and ROI before broader team commitment",
        "• Establish skill collection and agent development best practices"
    ]

    add_content_to_slide(slide15, content)

    # Save presentation
    output_path = r"c:\Users\adimi\OneDrive\Desktop\workspace\adept\Strategic_Analysis_Agentic_AI.pptx"
    prs.save(output_path)
    print("SUCCESS: Presentation created successfully: " + output_path)
    return output_path

def add_background(slide, color):
    """Add solid color background to slide"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_section_slide(prs, title):
    """Create a standard content slide with title"""
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)

    # Light background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_WHITE

    # Header bar
    header_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(10), Inches(1)
    )
    header_shape.fill.solid()
    header_shape.fill.fore_color.rgb = COLOR_PRIMARY
    header_shape.line.color.rgb = COLOR_PRIMARY

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.7))
    title_frame = title_box.text_frame
    title_frame.word_wrap = False
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE

    return slide

def add_content_to_slide(slide, content_list):
    """Add bullet point content to slide"""
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(6))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True

    for i, line in enumerate(content_list):
        if i > 0:
            p = text_frame.add_paragraph()
        else:
            p = text_frame.paragraphs[0]

        p.text = line

        # Determine indentation and formatting
        if line.startswith("•"):
            p.level = 1
            p.font.size = Pt(13)
        elif line.startswith("◦"):
            p.level = 2
            p.font.size = Pt(12)
        elif line == "":
            p.font.size = Pt(6)
        elif any(line.startswith(prefix) for prefix in ["Scenario", "Bottom Line", "Outcome", "Overall", "Recommended", "Why", "Alternative", "Strategic", "Skills", "Five", "Six", "Deliverables", "Phase"]):
            p.level = 0
            p.font.bold = True
            p.font.size = Pt(14)
            p.font.color.rgb = COLOR_PRIMARY
        else:
            p.level = 0
            p.font.size = Pt(13)

        p.font.color.rgb = COLOR_TEXT
        p.space_before = Pt(2)
        p.space_after = Pt(2)

if __name__ == "__main__":
    create_presentation()
